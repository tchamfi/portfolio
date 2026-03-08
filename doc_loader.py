"""
Auto-loader de documents pour le RAG "Ask Lionel".
Dépose tes fichiers (PDF, DOCX, PPTX, XLSX, TXT, MD) dans le dossier ./docs/
et ce script les découpe automatiquement en chunks pour ChromaDB.
"""

import os
import re
from pathlib import Path

# Dossier où déposer les documents
DOCS_FOLDER = "./docs"


def extract_text_from_pdf(filepath: str) -> str:
    """Extrait le texte d'un PDF."""
    try:
        from pypdf import PdfReader
    except ImportError:
        print("⚠️  Installation de pypdf...")
        os.system("pip install pypdf -q")
        from pypdf import PdfReader

    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def extract_text_from_docx(filepath: str) -> str:
    """Extrait le texte d'un fichier Word (.docx)."""
    try:
        from docx import Document
    except ImportError:
        print("⚠️  Installation de python-docx...")
        os.system("pip install python-docx -q")
        from docx import Document

    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Extrait aussi le texte des tableaux
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    return "\n".join(paragraphs)


def extract_text_from_pptx(filepath: str) -> str:
    """Extrait le texte d'une présentation PowerPoint (.pptx)."""
    try:
        from pptx import Presentation
    except ImportError:
        print("⚠️  Installation de python-pptx...")
        os.system("pip install python-pptx -q")
        from pptx import Presentation

    prs = Presentation(filepath)
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
            # Extrait le texte des tableaux dans les slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_parts.append(row_text)
        if len(slide_parts) > 1:  # Plus que juste le numéro de slide
            slides_text.append("\n".join(slide_parts))

    return "\n\n".join(slides_text)


def extract_text_from_xlsx(filepath: str) -> str:
    """Extrait le texte d'un fichier Excel (.xlsx)."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        print("⚠️  Installation de openpyxl...")
        os.system("pip install openpyxl -q")
        from openpyxl import load_workbook

    wb = load_workbook(filepath, data_only=True)
    sheets_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts = [f"[Feuille: {sheet_name}]"]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                sheet_parts.append(row_text)
        if len(sheet_parts) > 1:
            sheets_text.append("\n".join(sheet_parts))

    return "\n\n".join(sheets_text)


def extract_text_from_file(filepath: str) -> str:
    """Extrait le texte selon le type de fichier."""
    ext = Path(filepath).suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext == ".docx":
        return extract_text_from_docx(filepath)
    elif ext == ".pptx":
        return extract_text_from_pptx(filepath)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(filepath)
    elif ext in [".txt", ".md", ".csv", ".html", ".json", ".jsx"]:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    else:
        print(f"   ⏭️  Format ignoré : {ext}")
        return ""


def split_text_into_chunks(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    """
    Découpe un texte en chunks avec overlap.
    
    - chunk_size : nombre de caractères par chunk (~200 mots)
    - overlap : chevauchement entre chunks pour garder le contexte
    """
    # Nettoie le texte
    text = re.sub(r'\n{3,}', '\n\n', text)  # Max 2 sauts de ligne
    text = re.sub(r' {2,}', ' ', text)       # Max 1 espace

    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size

        # Essaie de couper à une fin de phrase ou de paragraphe
        if end < len(text):
            # Cherche le dernier point, retour à la ligne ou fin de phrase
            for sep in ['\n\n', '\n', '. ', '! ', '? ']:
                last_sep = text[start:end].rfind(sep)
                if last_sep > chunk_size * 0.5:  # Au moins à 50% du chunk
                    end = start + last_sep + len(sep)
                    break

        chunk_text = text[start:end].strip()
        if len(chunk_text) > 50:  # Ignore les chunks trop petits
            chunks.append(chunk_text)

        start = end - overlap  # Overlap pour garder le contexte

    return chunks


def load_documents_as_chunks() -> list[dict]:
    """
    Charge tous les documents du dossier ./docs/ et les convertit en chunks.
    Retourne une liste de dicts au même format que CV_CHUNKS.
    """
    # Crée le dossier s'il n'existe pas
    os.makedirs(DOCS_FOLDER, exist_ok=True)

    files = []
    for root, dirs, filenames in os.walk(DOCS_FOLDER):
        for f in filenames:
            if f.lower().endswith(('.pdf', '.docx', '.pptx', '.xlsx', '.txt', '.md', '.html', '.json', '.jsx')) and not f.startswith(('~', '.')):
                files.append(os.path.join(root, f))

    if not files:
        print(f"📂 Aucun document trouvé dans {DOCS_FOLDER}/")
        return []

    all_chunks = []

    for filepath in sorted(files):
        # Chemin relatif pour l'affichage
        rel_path = os.path.relpath(filepath, DOCS_FOLDER)
        filename = os.path.basename(filepath)
        print(f"📄 Lecture de : {rel_path}")

        # Extrait le texte
        text = extract_text_from_file(filepath)
        if not text:
            continue

        # Découpe en chunks
        text_chunks = split_text_into_chunks(text)
        print(f"   → {len(text_chunks)} chunks créés")

        # Génère un nom propre pour la catégorie
        doc_name = Path(filename).stem.replace("_", " ").replace("-", " ")

        for i, chunk_text in enumerate(text_chunks):
            safe_name = rel_path.replace("/", "_").replace(" ", "_").replace(".", "_")
            chunk_id = f"doc_{safe_name}_{i:03d}"
            all_chunks.append({
                "id": chunk_id,
                "text": chunk_text,
                "metadata": {
                    "category": "document",
                    "source": rel_path,
                    "doc_name": doc_name,
                    "chunk_index": i,
                    "keywords": doc_name,
                }
            })

    print(f"\n✅ Total : {len(all_chunks)} chunks chargés depuis {len(files)} fichier(s)")
    return all_chunks


# --- Test en standalone ---
if __name__ == "__main__":
    chunks = load_documents_as_chunks()
    for c in chunks[:3]:  # Affiche les 3 premiers chunks
        print(f"\n--- {c['id']} ({c['metadata']['source']}) ---")
        print(c['text'][:200] + "...")
